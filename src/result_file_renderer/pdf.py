"""
Pipeline из 3 шагов:
  1) читает локальный JSON с ответами;
  2) записывает ответы в Excel-методику, столбец D,
     считает метрики и top-5 по рейтингу из столбца I;
  3) вставляет метрики в PPTX-шаблон и экспортирует PPTX в PDF.

Входной JSON:
{
  "b1_1": {
    "question_text": "...",
    "user_answer": "2"
  }
}

Перед запуском установить зависимости:

pip install openpyxl python-pptx pillow

Для экспорта PDF нужен LibreOffice/soffice в PATH.
"""

import json
import math
import re
import shutil
import subprocess
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

from openpyxl import load_workbook
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

FIRST_DATA_ROW = 2
ANSWER_MIN = 0.0
ANSWER_MAX = 3.0
DEFAULT_STUDY_URL = "https://yandex.ru"


@dataclass(frozen=True)
class SkillMetric:
    rank: int | None
    row_number: int
    skill_id: int
    domain: str
    skill: str
    answer: float
    answer_norm: float
    inverse_answer_norm: float
    current_demand: float
    dsq_norm: float
    rating: float


def load_answers(json_path: str | Path) -> dict[str, Any]:
    json_path = Path(json_path)
    with json_path.open("r", encoding="utf-8") as file:
        data = json.load(file)

    if not isinstance(data, dict):
        raise ValueError(
            "JSON должен быть объектом формата {question_id: {question_text, user_answer}}"
        )
    return data


def normalize_text(text: str) -> str:
    text = text.lower().replace("ё", "е")
    return re.sub(r"[^a-zа-я0-9]+", " ", text).strip()


def to_float(value: Any, field_name: str) -> float:
    if value is None or value == "":
        raise ValueError(f"Пустое значение в поле {field_name}")
    if isinstance(value, str):
        value = value.replace(",", ".").strip()
    try:
        return float(value)
    except (TypeError, ValueError) as exc:
        raise ValueError(
            f"Не удалось прочитать число в поле {field_name}: {value!r}"
        ) from exc


def get_skill_rows(ws: Worksheet) -> list[int]:
    rows: list[int] = []
    for row in range(FIRST_DATA_ROW, ws.max_row + 1):
        value = ws.cell(row=row, column=1).value
        if isinstance(value, int):
            rows.append(row)
        elif isinstance(value, float) and value.is_integer():
            rows.append(row)
        elif isinstance(value, str) and value.strip().isdigit():
            rows.append(row)

    if not rows:
        raise ValueError(
            "В Excel не найдены строки навыков: ожидались номера в столбце A"
        )
    return rows


def build_question_key_to_row(ws: Worksheet, skill_rows: list[int]) -> dict[str, int]:
    domains: list[str] = []
    for row in skill_rows:
        domain = str(ws.cell(row=row, column=2).value).strip()
        if domain not in domains:
            domains.append(domain)

    mapping: dict[str, int] = {}
    for domain_idx, domain in enumerate(domains, start=1):
        local_idx = 0
        for row in skill_rows:
            if str(ws.cell(row=row, column=2).value).strip() == domain:
                local_idx += 1
                mapping[f"b{domain_idx}_{local_idx}"] = row
    return mapping


def build_question_text_to_row(ws: Worksheet, skill_rows: list[int]) -> dict[str, int]:
    mapping: dict[str, int] = {}
    for row in skill_rows:
        skill = str(ws.cell(row=row, column=3).value).strip()
        if skill:
            mapping[normalize_text(skill)] = row
    return mapping


def resolve_row_for_answer(
    question_key: str,
    payload: Any,
    key_to_row: dict[str, int],
    text_to_row: dict[str, int],
    skill_id_to_row: dict[int, int],
) -> int | None:
    """Сопоставляет ответ из JSON со строкой Excel."""
    if question_key in key_to_row:
        return key_to_row[question_key]

    match = re.search(r"(?:^|_)(\d+)$", question_key)
    if match:
        skill_id = int(match.group(1))
        if skill_id in skill_id_to_row:
            return skill_id_to_row[skill_id]

    if isinstance(payload, dict):
        question_text = payload.get("question_text")
        if isinstance(question_text, str):
            return text_to_row.get(normalize_text(question_text))

    return None


def extract_answer(payload: Any, question_key: str) -> float:
    raw_value = payload.get("user_answer") if isinstance(payload, dict) else payload
    answer = to_float(raw_value, f"{question_key}.user_answer")
    if answer < ANSWER_MIN or answer > ANSWER_MAX:
        raise ValueError(
            f"Ответ {question_key} должен быть в диапазоне {ANSWER_MIN:g}..{ANSWER_MAX:g}, получено: {answer}"
        )
    return answer


def write_answers_to_column_d(
    ws: Worksheet,
    answers: dict[str, Any],
) -> tuple[list[int], set[int], list[str]]:
    skill_rows = get_skill_rows(ws)
    key_to_row = build_question_key_to_row(ws, skill_rows)
    text_to_row = build_question_text_to_row(ws, skill_rows)
    skill_id_to_row = {int(ws.cell(row=row, column=1).value): row for row in skill_rows}

    matched_rows: set[int] = set()
    warnings: list[str] = []

    for question_key, payload in answers.items():
        row = resolve_row_for_answer(
            question_key, payload, key_to_row, text_to_row, skill_id_to_row
        )
        if row is None:
            warnings.append(
                f"Ответ {question_key!r} не сопоставлен со строкой Excel и пропущен"
            )
            continue

        ws.cell(row=row, column=4).value = extract_answer(payload, question_key)
        matched_rows.add(row)

    missing_rows = [row for row in skill_rows if row not in matched_rows]
    if missing_rows:
        missing_ids = [ws.cell(row=row, column=1).value for row in missing_rows]
        raise ValueError(f"Нет ответов для навыков: {missing_ids}")

    for row in missing_rows:
        ws.cell(row=row, column=4).value = 0.0

    if missing_rows:
        missing_ids = [str(ws.cell(row=row, column=1).value) for row in missing_rows]
        warnings.append(
            "Для отсутствующих навыков подставлено 0: " + ", ".join(missing_ids)
        )

    return skill_rows, matched_rows, warnings


def ensure_excel_formulas(ws: Worksheet, skill_rows: list[int]) -> None:
    for row in skill_rows:
        ws.cell(row=row, column=5).value = f"=D{row}/3"
        ws.cell(row=row, column=6).value = f"=1-E{row}"
        ws.cell(row=row, column=9).value = (
            f"=IF(COUNTIF(F{row}:H{row},0)>0,0,GEOMEAN(F{row}:H{row}))"
        )


def geomean_or_zero(values: list[float]) -> float:
    if any(value <= 0 for value in values):
        return 0.0
    return math.prod(values) ** (1 / len(values))


def calculate_metrics(
    ws: Worksheet,
    skill_rows: list[int],
    *,
    matched_rows: set[int],
) -> dict[str, Any]:
    skill_metrics: list[SkillMetric] = []

    for row in skill_rows:
        skill_id = int(ws.cell(row=row, column=1).value)
        domain = str(ws.cell(row=row, column=2).value).strip()
        skill = str(ws.cell(row=row, column=3).value).strip()
        answer = to_float(ws.cell(row=row, column=4).value, f"D{row}")
        current_demand = to_float(ws.cell(row=row, column=7).value, f"G{row}")
        dsq_norm = to_float(ws.cell(row=row, column=8).value, f"H{row}")

        answer_norm = answer / 3.0
        inverse_answer_norm = 1.0 - answer_norm
        rating = geomean_or_zero([inverse_answer_norm, current_demand, dsq_norm])

        skill_metrics.append(
            SkillMetric(
                rank=None,
                row_number=row,
                skill_id=skill_id,
                domain=domain,
                skill=skill,
                answer=answer,
                answer_norm=answer_norm,
                inverse_answer_norm=inverse_answer_norm,
                current_demand=current_demand,
                dsq_norm=dsq_norm,
                rating=rating,
            )
        )
    top5_candidates = skill_metrics

    if len(top5_candidates) < 5:
        raise ValueError(
            f"Недостаточно строк для top-5: найдено {len(top5_candidates)}. "
            "Передайте больше ответов."
        )

    top5_metrics = sorted(
        top5_candidates, key=lambda item: (-item.rating, item.skill_id)
    )[:5]
    top5: list[SkillMetric] = []
    for rank, item in enumerate(top5_metrics, start=1):
        top5.append(
            SkillMetric(
                rank=rank,
                row_number=item.row_number,
                skill_id=item.skill_id,
                domain=item.domain,
                skill=item.skill,
                answer=item.answer,
                answer_norm=item.answer_norm,
                inverse_answer_norm=item.inverse_answer_norm,
                current_demand=item.current_demand,
                dsq_norm=item.dsq_norm,
                rating=item.rating,
            )
        )

    domains: list[str] = []
    for item in skill_metrics:
        if item.domain not in domains:
            domains.append(item.domain)

    domain_averages = {
        domain: sum(item.answer for item in skill_metrics if item.domain == domain)
        / len([item for item in skill_metrics if item.domain == domain])
        for domain in domains
    }
    overall_average = sum(item.answer for item in skill_metrics) / len(skill_metrics)

    return {
        "overall_average": overall_average,
        "overall_average_formula": f"SUM(D)/{len(skill_metrics)}",
        "skill_count": len(skill_metrics),
        "matched_answer_count": len(matched_rows),
        "domain_averages": domain_averages,
        "top5": [asdict(item) for item in top5],
        "skills": [asdict(item) for item in skill_metrics],
    }


def recreate_metrics_sheet(wb: Any, metrics: dict[str, Any]) -> None:
    sheet_name = "Метрики"
    if sheet_name in wb.sheetnames:
        del wb[sheet_name]
    ws = wb.create_sheet(sheet_name)

    ws["A1"] = "Итоговый средний результат"
    ws["B1"] = metrics["overall_average"]
    ws["B1"].number_format = "0.00"

    ws["A3"] = "Средние по доменам"
    ws["A4"] = "Домен"
    ws["B4"] = "Средний ответ"
    row = 5
    for domain, value in metrics["domain_averages"].items():
        ws.cell(row=row, column=1).value = domain
        ws.cell(row=row, column=2).value = value
        ws.cell(row=row, column=2).number_format = "0.00"
        row += 1

    top_start = row + 2
    scope_label = (
        "все строки" if metrics.get("top5_scope") == "all" else "только строки из JSON"
    )
    ws.cell(row=top_start, column=1).value = (
        f"Top-5 по рейтингу из столбца I ({scope_label})"
    )
    headers = ["Ранг", "№ навыка", "Домен", "Навык", "Ответ D", "Рейтинг I"]
    for col, header in enumerate(headers, start=1):
        ws.cell(row=top_start + 1, column=col).value = header

    for idx, item in enumerate(metrics["top5"], start=top_start + 2):
        ws.cell(row=idx, column=1).value = item["rank"]
        ws.cell(row=idx, column=2).value = item["skill_id"]
        ws.cell(row=idx, column=3).value = item["domain"]
        ws.cell(row=idx, column=4).value = item["skill"]
        ws.cell(row=idx, column=5).value = item["answer"]
        ws.cell(row=idx, column=6).value = item["rating"]
        ws.cell(row=idx, column=6).number_format = "0.000000"

    widths = {"A": 22, "B": 16, "C": 34, "D": 78, "E": 12, "F": 14}
    for column, width in widths.items():
        ws.column_dimensions[column].width = width


def force_excel_recalculate_on_open(wb: Any) -> None:
    """Просит Excel/LibreOffice пересчитать формулы при открытии файла."""
    try:
        wb.calculation.fullCalcOnLoad = True
        wb.calculation.forceFullCalc = True
        wb.calculation.calcMode = "auto"
    except Exception:
        pass


def fill_excel_and_calculate_metrics(
    answers_path: str | Path,
    methodology_path: str | Path,
    out_xlsx: str | Path,
    out_metrics: str | Path,
    *,
    sheet_name: str | None,
) -> dict[str, Any]:
    answers = load_answers(answers_path)
    wb = load_workbook(methodology_path)
    ws = wb[sheet_name] if sheet_name else wb.active

    skill_rows, matched_rows, warnings = write_answers_to_column_d(ws, answers)
    ensure_excel_formulas(ws, skill_rows)
    metrics = calculate_metrics(
        ws,
        skill_rows,
        matched_rows=matched_rows,
    )
    if warnings:
        metrics["warnings"] = warnings

    recreate_metrics_sheet(wb, metrics)
    force_excel_recalculate_on_open(wb)

    out_xlsx = Path(out_xlsx)
    out_metrics = Path(out_metrics)
    out_xlsx.parent.mkdir(parents=True, exist_ok=True)
    out_metrics.parent.mkdir(parents=True, exist_ok=True)

    wb.save(out_xlsx)
    with out_metrics.open("w", encoding="utf-8") as file:
        json.dump(metrics, file, ensure_ascii=False, indent=2)

    return metrics


def format_ru_number(value: float, digits: int = 1) -> str:
    return f"{float(value):.{digits}f}".replace(".", ",")


def replace_paragraph_text(paragraph: Any, text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = text


def set_shape_text(shape: Any, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    paragraphs = shape.text_frame.paragraphs
    if not paragraphs:
        return
    replace_paragraph_text(paragraphs[0], text)
    for paragraph in paragraphs[1:]:
        replace_paragraph_text(paragraph, "")


def clear_paragraph(paragraph: Any) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = ""
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = ""


def add_link_run(
    paragraph: Any, text: str, url: str, *, like_run: Any | None = None
) -> Any:
    run = paragraph.add_run()
    run.text = text
    run.hyperlink.address = url

    if like_run is not None:
        for attr in ("size", "name", "bold", "italic"):
            try:
                setattr(run.font, attr, getattr(like_run.font, attr))
            except Exception:
                pass

    try:
        run.font.color.rgb = RGBColor(5, 99, 193)
        run.font.underline = True
    except Exception:
        pass
    return run


def set_research_link_in_intro(slide: Any, url: str) -> None:
    shape = next(
        (
            s
            for s in slide.shapes
            if getattr(s, "has_text_frame", False)
            and "Ссылка на исследование" in s.text
        ),
        None,
    )
    if shape is None:
        return

    paragraph = shape.text_frame.paragraphs[0]
    clear_paragraph(paragraph)
    prefix = (
        "*Весной 2026 года мы провели исследование среди руководителей российских компаний, "
        "чтобы оценить их восприятие уровня владения управленческими навыками, их текущей "
        "и будущей востребованности. "
    )
    if paragraph.runs:
        paragraph.runs[0].text = prefix
    else:
        paragraph.add_run().text = prefix
    add_link_run(paragraph, "Ссылка на исследование", url, like_run=paragraph.runs[0])


def set_result_text(slide: Any, overall_average: float) -> None:
    value = format_ru_number(overall_average, 2)
    candidates = [
        s
        for s in slide.shapes
        if getattr(s, "has_text_frame", False) and s.text.strip() == "2,0"
    ]
    if not candidates:
        # fallback: ищем текстовое поле, похожее на число
        candidates = [
            s
            for s in slide.shapes
            if getattr(s, "has_text_frame", False)
            and s.text.strip().replace(",", ".").replace(".", "", 1).isdigit()
        ]
    if candidates:
        shape = candidates[0]
        try:
            shape.width = int(shape.width * 1.10)
            shape.text_frame.word_wrap = False
            shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        set_shape_text(shape, value)
        try:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(58)
        except Exception:
            pass


def get_template_chart_data(slide: Any) -> tuple[list[str], list[float], str, str]:
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            chart = shape.chart
            categories = [str(c) for c in chart.plots[0].categories]
            series = list(chart.series)
            if len(series) >= 2:
                return (
                    categories,
                    [float(v) for v in series[0].values],
                    str(series[0].name),
                    str(series[1].name),
                )
    raise ValueError("В шаблоне не найдена диаграмма профиля навыков")


def update_profile_chart(slide: Any, metrics: dict[str, Any]) -> None:
    chart_shape = next(
        (s for s in slide.shapes if getattr(s, "has_chart", False)), None
    )
    if chart_shape is None:
        raise ValueError("В шаблоне не найдена диаграмма")

    categories, survey_values, survey_name, individual_name = get_template_chart_data(
        slide
    )
    domain_averages = metrics["domain_averages"]
    individual_values = [
        round(float(domain_averages.get(category, 0.0)), 2) for category in categories
    ]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(survey_name, survey_values)
    chart_data.add_series(individual_name, individual_values)

    chart = chart_shape.chart
    chart.replace_data(chart_data)

    try:
        chart.value_axis.minimum_scale = 0.0
        chart.value_axis.maximum_scale = 3.0
        chart.value_axis.major_unit = 1.0
    except Exception:
        pass


def update_gauge_needle(slide: Any, overall_average: float) -> None:
    needle_candidates = [s for s in slide.shapes if "треугольник" in s.name.lower()]
    if not needle_candidates:
        return
    clamped = max(0.0, min(3.0, float(overall_average)))
    # Шкала 0..3: 0 -> -90 градусов, 1.5 -> 0 градусов, 3 -> 90 градусов.
    needle_candidates[0].rotation = clamped / 3.0 * 180.0 - 90.0


def set_top5_text(slide: Any, top5: list[dict[str, Any]], url: str) -> None:
    shape = next(
        (
            s
            for s in slide.shapes
            if getattr(s, "has_text_frame", False)
            and "Какие навыки развивать" in s.text
        ),
        None,
    )
    if shape is None:
        raise ValueError("В шаблоне не найден текстовый блок top-5")

    tf = shape.text_frame
    while len(tf.paragraphs) < 9:
        tf.add_paragraph()

    replace_paragraph_text(tf.paragraphs[0], "Какие навыки развивать в первую очередь?")
    for i, item in enumerate(top5[:5], start=1):
        replace_paragraph_text(tf.paragraphs[i], str(item["skill"]))

    replace_paragraph_text(tf.paragraphs[6], "")
    replace_paragraph_text(tf.paragraphs[7], "")

    # В последнем абзаце делаем слово 'исследовании' активной ссылкой.
    note_p = tf.paragraphs[8]
    clear_paragraph(note_p)
    note_prefix = (
        "Рассчитано на основе ваших личных показателей и средних значений текущей "
        "и ожидаемой востребованности навыка. Подробности можно узнать в нашем "
    )
    if note_p.runs:
        note_p.runs[0].text = note_prefix
    else:
        note_p.add_run().text = note_prefix
    add_link_run(note_p, "исследовании", url, like_run=note_p.runs[0])
    note_p.add_run().text = "."

    for paragraph in tf.paragraphs[9:]:
        replace_paragraph_text(paragraph, "")


def add_child(
    parent: Any, ns: str, tag: str, attrs: dict[str, str] | None = None
) -> Any:
    """Добавляет XML-элемент в нужном namespace."""
    from xml.etree import ElementTree as ET

    child = ET.SubElement(parent, f"{{{ns}}}{tag}")
    if attrs:
        child.attrib.update(attrs)
    return child


def build_value_data_labels_xml(chart_ns: str, drawing_ns: str) -> Any:
    from xml.etree import ElementTree as ET

    d_lbls = ET.Element(f"{{{chart_ns}}}dLbls")
    add_child(
        d_lbls,
        chart_ns,
        "numFmt",
        {"formatCode": "0.00", "sourceLinked": "0"},
    )

    sp_pr = add_child(d_lbls, chart_ns, "spPr")
    add_child(sp_pr, drawing_ns, "noFill")
    line = add_child(sp_pr, drawing_ns, "ln")
    add_child(line, drawing_ns, "noFill")
    add_child(sp_pr, drawing_ns, "effectLst")

    tx_pr = add_child(d_lbls, chart_ns, "txPr")
    add_child(
        tx_pr,
        drawing_ns,
        "bodyPr",
        {
            "rot": "0",
            "spcFirstLastPara": "1",
            "vertOverflow": "ellipsis",
            "vert": "horz",
            "wrap": "square",
            "lIns": "38100",
            "tIns": "19050",
            "rIns": "38100",
            "bIns": "19050",
            "anchor": "ctr",
            "anchorCtr": "1",
        },
    )
    add_child(tx_pr, drawing_ns, "lstStyle")
    paragraph = add_child(tx_pr, drawing_ns, "p")
    p_pr = add_child(paragraph, drawing_ns, "pPr")
    def_r_pr = add_child(
        p_pr,
        drawing_ns,
        "defRPr",
        {
            "sz": "1000",
            "b": "0",
            "i": "0",
            "u": "none",
            "strike": "noStrike",
            "kern": "1200",
            "baseline": "0",
        },
    )
    solid_fill = add_child(def_r_pr, drawing_ns, "solidFill")
    add_child(solid_fill, drawing_ns, "srgbClr", {"val": "404040"})
    add_child(def_r_pr, drawing_ns, "latin", {"typeface": "+mn-lt"})
    add_child(def_r_pr, drawing_ns, "ea", {"typeface": "+mn-ea"})
    add_child(def_r_pr, drawing_ns, "cs", {"typeface": "+mn-cs"})
    add_child(paragraph, drawing_ns, "endParaRPr", {"lang": "ru-RU"})

    for name, value in (
        ("showLegendKey", "0"),
        ("showVal", "1"),
        ("showCatName", "0"),
        ("showSerName", "0"),
        ("showPercent", "0"),
        ("showBubbleSize", "0"),
        ("showLeaderLines", "0"),
    ):
        add_child(d_lbls, chart_ns, name, {"val": value})

    return d_lbls


def remove_existing_data_labels(root: Any, chart_ns: str) -> None:
    d_lbls_tag = f"{{{chart_ns}}}dLbls"
    for parent in root.iter():
        for child in list(parent):
            if child.tag == d_lbls_tag:
                parent.remove(child)


def insert_series_data_labels(series: Any, d_lbls: Any, chart_ns: str) -> None:
    marker_tag = f"{{{chart_ns}}}marker"
    cat_tag = f"{{{chart_ns}}}cat"

    children = list(series)
    insert_index = None

    for idx, child in enumerate(children):
        if child.tag == marker_tag:
            insert_index = idx + 1
            break

    if insert_index is None:
        for idx, child in enumerate(children):
            if child.tag == cat_tag:
                insert_index = idx
                break

    if insert_index is None:
        series.append(d_lbls)
    else:
        series.insert(insert_index, d_lbls)


def configure_profile_chart_data_labels(pptx_path: str | Path) -> None:
    pptx_path = Path(pptx_path)
    tmp_path = pptx_path.with_suffix(".tmp.pptx")
    chart_ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(
        tmp_path, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/charts/") and item.filename.endswith(
                ".xml"
            ):
                try:
                    from xml.etree import ElementTree as ET

                    ET.register_namespace("c", chart_ns)
                    ET.register_namespace("a", drawing_ns)
                    root = ET.fromstring(data)
                    remove_existing_data_labels(root, chart_ns)

                    radar_charts = root.findall(f".//{{{chart_ns}}}radarChart")
                    for radar_chart in radar_charts:
                        series_list = radar_chart.findall(f"{{{chart_ns}}}ser")
                        if len(series_list) >= 2:
                            individual_series = series_list[1]
                            d_lbls = build_value_data_labels_xml(chart_ns, drawing_ns)
                            insert_series_data_labels(
                                individual_series, d_lbls, chart_ns
                            )

                    data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(pptx_path)


def fill_template(
    template_pptx: str | Path,
    metrics: dict[str, Any],
    out_pptx: str | Path,
    study_url: str,
) -> None:
    """Шаг 3A: вставляет метрики в PPTX-шаблон."""
    prs = Presentation(template_pptx)
    if not prs.slides:
        raise ValueError("В PPTX-шаблоне нет слайдов")

    slide = prs.slides[0]
    set_result_text(slide, metrics["overall_average"])
    update_profile_chart(slide, metrics)
    update_gauge_needle(slide, metrics["overall_average"])
    set_top5_text(slide, metrics["top5"], study_url)
    set_research_link_in_intro(slide, study_url)

    out_pptx = Path(out_pptx)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)
    configure_profile_chart_data_labels(out_pptx)


def find_office_binary() -> str:
    for binary in ("libreoffice", "soffice"):
        path = shutil.which(binary)
        if path:
            return path
    raise RuntimeError("Не найден LibreOffice/soffice")


def convert_pptx_to_pdf(pptx_path: str | Path, pdf_path: str | Path) -> None:
    pptx_path = Path(pptx_path).resolve()
    pdf_path = Path(pdf_path).resolve()
    pdf_path.parent.mkdir(parents=True, exist_ok=True)

    office = find_office_binary()
    cmd = [
        office,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_path.parent),
        str(pptx_path),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(
            "Ошибка экспорта PPTX в PDF:\n"
            f"STDOUT:\n{result.stdout}\nSTDERR:\n{result.stderr}"
        )

    generated_pdf = pdf_path.parent / f"{pptx_path.stem}.pdf"
    if generated_pdf != pdf_path:
        if pdf_path.exists():
            pdf_path.unlink()
        generated_pdf.rename(pdf_path)

    if not pdf_path.exists():
        raise RuntimeError("LibreOffice завершился без ошибки, но PDF-файл не найден")


@dataclass(frozen=True)
class ReportResult:
    """Результат выполнения полного pipeline."""

    xlsx_path: str
    metrics_path: str
    pptx_path: str
    pdf_path: str | None
    study_url: str
    top5_scope: str
    overall_average: float
    warnings: list[str]
    metrics: dict[str, Any]


def generate_management_index_report(
    *,
    answers_path: str | Path,
    methodology_path: str | Path,
    template_pptx_path: str | Path,
    output_dir: str | Path,
    base_name: str = "management_index_report",
    sheet_name: str | None = None,
    study_url: str = DEFAULT_STUDY_URL,
    export_pdf: bool = True,
) -> ReportResult:

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    out_xlsx = output_dir / f"{base_name}_filled_methodology.xlsx"
    out_metrics = output_dir / f"{base_name}_metrics.json"
    out_pptx = output_dir / f"{base_name}.pptx"
    out_pdf = output_dir / f"{base_name}.pdf"

    print("Шаг 1/3: читаю JSON с ответами")
    print("Шаг 2/3: заполняю Excel, столбец D, и считаю метрики")
    metrics = fill_excel_and_calculate_metrics(
        answers_path=answers_path,
        methodology_path=methodology_path,
        out_xlsx=out_xlsx,
        out_metrics=out_metrics,
        sheet_name=sheet_name,
    )

    print("Шаг 3/3: заполняю PPTX-шаблон")
    fill_template(template_pptx_path, metrics, out_pptx, study_url)

    pdf_path: str | None = None
    if export_pdf:
        convert_pptx_to_pdf(out_pptx, out_pdf)
        pdf_path = str(out_pdf)

    return ReportResult(
        xlsx_path=str(out_xlsx),
        metrics_path=str(out_metrics),
        pptx_path=str(out_pptx),
        pdf_path=pdf_path,
        study_url=study_url,
        top5_scope=str(metrics.get("top5_scope")),
        overall_average=float(metrics["overall_average"]),
        warnings=list(metrics.get("warnings", [])),
        metrics=metrics,
    )


def print_report_summary(result: ReportResult) -> None:
    """Печатает краткую сводку для локального запуска."""
    print("Готово:")
    print(f"- Excel:   {result.xlsx_path}")
    print(f"- Metrics: {result.metrics_path}")
    print(f"- PPTX:    {result.pptx_path}")
    if result.pdf_path:
        print(f"- PDF:     {result.pdf_path}")
    print(f"- Link:    {result.study_url}")
    print(f"- Top-5 scope: {result.top5_scope}")
    print(f"- Итоговая оценка: {result.overall_average:.2f}")
    print("\nTop-5 по рейтингу I:")
    for item in result.metrics["top5"]:
        print(
            f"{item['rank']}. №{item['skill_id']} — {item['skill']} — I={item['rating']:.6f}"
        )

    if result.warnings:
        print("\nПредупреждения:")
        for warning in result.warnings:
            print(f"- {warning}")

    return 0


if __name__ == "__main__":
    ANSWERS_PATH = Path("result_file_renderer/example_answers.json")
    METHODOLOGY_PATH = Path("result_file_renderer/template.xlsx")
    TEMPLATE_PPTX_PATH = Path("result_file_renderer/template.pptx")
    OUTPUT_DIR = Path("result_file_renderer/results")
    BASE_NAME = "report"
    STUDY_URL = "https://yandex.ru"
    EXPORT_PDF = True
    SHEET_NAME = None

    result = generate_management_index_report(
        answers_path=ANSWERS_PATH,
        methodology_path=METHODOLOGY_PATH,
        template_pptx_path=TEMPLATE_PPTX_PATH,
        output_dir=OUTPUT_DIR,
        base_name=BASE_NAME,
        sheet_name=SHEET_NAME,
        study_url=STUDY_URL,
        export_pdf=EXPORT_PDF,
    )
    print_report_summary(result)
